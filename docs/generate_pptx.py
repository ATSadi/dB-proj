"""Generate the Campus Maintenance project PPTX presentation."""
from __future__ import annotations

import os

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = r"d:\dB proj"
ASSETS = os.path.join(ROOT, "docs", "report_assets")
LABELED = os.path.join(ASSETS, "labeled")
OUT_PPTX = os.path.join(
    ROOT, "docs", "Campus_Maintenance_Presentation_Ahnaf_Tajwar_Sadi_2207104.pptx"
)
KUET = os.path.join(ASSETS, "kuet_logo.png")
GITHUB = "https://github.com/ATSadi/dB-proj/"

BG = RGBColor(0x0B, 0x12, 0x20)
CARD = RGBColor(0x1E, 0x29, 0x3B)
TEXT = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
LINE = RGBColor(0x33, 0x41, 0x55)


def labeled(name: str) -> str:
    return os.path.join(LABELED, name)


def asset(name: str) -> str:
    return os.path.join(ASSETS, name)


def set_run(run, size=18, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return box


def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    return shape


def add_pic(slide, src, left, top, max_w, max_h):
    if not os.path.exists(src):
        print("missing", src)
        return None
    with PILImage.open(src) as im:
        w_px, h_px = im.size
    max_w_in = max_w.inches
    max_h_in = max_h.inches
    scale = min(max_w_in / (w_px / 96.0), max_h_in / (h_px / 96.0))
    return slide.shapes.add_picture(
        src,
        left,
        top,
        width=Inches((w_px / 96.0) * scale),
        height=Inches((h_px / 96.0) * scale),
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        add_bg(s, prs)
        return s

    # 1 Title
    s = slide()
    if os.path.exists(KUET):
        s.shapes.add_picture(KUET, Inches(6.15), Inches(0.35), height=Inches(1.05))
    add_text(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.4), "Good Morning", 20, False, ACCENT, PP_ALIGN.CENTER)
    add_text(
        s,
        Inches(0.8),
        Inches(2.2),
        Inches(11.7),
        Inches(1.3),
        "Campus Maintenance &\nComplaint Management System",
        34,
        True,
        TEXT,
        PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(3.75),
        Inches(11.7),
        Inches(0.4),
        "Database Systems Lab Project  ·  Oracle SQL/PL/SQL  ·  Node.js  ·  Web Portals",
        15,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.15),
        Inches(11.7),
        Inches(0.9),
        "Name: Ahnaf Tajwar Sadi\nRoll: 2207104  ·  KUET CSE",
        18,
        True,
        TEXT,
        PP_ALIGN.CENTER,
    )
    add_text(s, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.35), GITHUB, 12, False, PURPLE, PP_ALIGN.CENTER)

    # 2 Overview
    s = slide()
    add_text(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.55), "Project Overview", 30, True)
    add_text(
        s,
        Inches(0.7),
        Inches(1.1),
        Inches(12),
        Inches(0.7),
        "A database-driven campus system for complaints, worker assignments, SLA tracking, chronic-issue detection, and monthly reporting.",
        17,
        False,
        MUTED,
    )
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.1),
        Inches(11.5),
        Inches(4.8),
        [
            "Students report location-based maintenance issues with priority",
            "Supervisors assign workers and monitor overdue / chronic cases",
            "Workers execute jobs and update availability",
            "Students rate completed work; Admins generate monthly reports",
            "Oracle triggers/procedures enforce SLA, audit, scoring, and escalation",
        ],
        20,
    )

    # 3 Features
    s = slide()
    add_text(s, Inches(0.55), Inches(0.3), Inches(12), Inches(0.55), "Key Features", 30, True)
    cards = [
        (0.45, "Authentication", ["Email + password login", "Forgot-password reset", "Role-based session guards"]),
        (3.55, "Student Portal", ["Submit complaints", "Filter / search / SLA", "Rate resolved jobs"]),
        (6.65, "Worker Portal", ["Active & done queues", "Start / resolve jobs", "Availability toggle"]),
        (9.75, "Ops + Admin", ["Assign / overdue / chronic", "Monthly reports", "Directory + CSV import"]),
    ]
    for left, title, items in cards:
        add_card(s, Inches(left), Inches(1.15), Inches(2.95), Inches(5.6))
        add_text(s, Inches(left + 0.18), Inches(1.4), Inches(2.6), Inches(0.55), title, 17, True, ACCENT)
        add_bullets(s, Inches(left + 0.18), Inches(2.2), Inches(2.6), Inches(4.2), items, 14)

    # 4 Workflow
    s = slide()
    add_text(s, Inches(0.5), Inches(0.22), Inches(12), Inches(0.45), "End-to-End Complaint Lifecycle", 26, True)
    add_pic(s, asset("diagram-workflow.png"), Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.1))

    # 5 Roles
    s = slide()
    add_text(s, Inches(0.5), Inches(0.22), Inches(12), Inches(0.45), "Role-Based Access Model", 26, True)
    add_pic(s, asset("diagram-roles.png"), Inches(0.6), Inches(0.9), Inches(12.1), Inches(6.1))

    # 6 Architecture
    s = slide()
    add_text(s, Inches(0.5), Inches(0.22), Inches(12), Inches(0.45), "System Architecture", 26, True)
    add_pic(s, asset("diagram-architecture.png"), Inches(0.7), Inches(0.9), Inches(12.0), Inches(6.1))

    # 7 ER overview
    s = slide()
    add_text(s, Inches(0.45), Inches(0.18), Inches(12), Inches(0.4), "Entity Relationship Overview", 24, True)
    add_pic(s, asset("diagram-erd-overview.png"), Inches(0.35), Inches(0.7), Inches(12.6), Inches(6.4))

    # 8 Schema diagram
    s = slide()
    add_text(s, Inches(0.45), Inches(0.18), Inches(12), Inches(0.4), "Database Schema Diagram", 24, True)
    add_text(
        s,
        Inches(0.5),
        Inches(0.65),
        Inches(12),
        Inches(0.35),
        "USERS · LOCATIONS · COMPLAINTS · WORKERS · ASSIGNMENTS · STATUS_LOG · FEEDBACK · CHRONIC_FLAGS · MAINTENANCE_REPORTS",
        12,
        False,
        MUTED,
    )
    add_pic(s, labeled("diagram-erd-dbeaver.png"), Inches(1.0), Inches(1.1), Inches(11.3), Inches(5.9))

    # 9 PL/SQL
    s = slide()
    add_text(s, Inches(0.5), Inches(0.22), Inches(12), Inches(0.45), "PL/SQL Business Logic", 26, True)
    add_pic(s, asset("diagram-plsql.png"), Inches(0.55), Inches(0.85), Inches(12.2), Inches(6.2))

    # 10 Login
    s = slide()
    add_text(s, Inches(0.5), Inches(0.22), Inches(12), Inches(0.45), "Secure Login & Demo Access", 26, True)
    add_text(
        s,
        Inches(0.55),
        Inches(0.75),
        Inches(12),
        Inches(0.35),
        "Email + password authentication  ·  Default demo password: Password123  ·  Forgot-password reset codes",
        13,
        False,
        MUTED,
    )
    add_pic(s, labeled("ui-01-login.png"), Inches(1.5), Inches(1.25), Inches(10.3), Inches(5.7))

    # 11 Student
    s = slide()
    add_text(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.4), "Student Portal", 26, True)
    add_pic(s, labeled("ui-02-student.png"), Inches(0.35), Inches(0.7), Inches(6.25), Inches(3.2))
    add_pic(s, labeled("ui-04-student-detail.png"), Inches(6.8), Inches(0.7), Inches(6.15), Inches(3.2))
    add_bullets(
        s,
        Inches(0.7),
        Inches(4.2),
        Inches(11.8),
        Inches(2.8),
        [
            "Submit complaints tied to campus locations with category and priority",
            "Filter, search, and track SLA remaining time on open tickets",
            "Rate resolved work to feed worker performance scoring",
        ],
        18,
    )

    # 12 Worker
    s = slide()
    add_text(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.4), "Worker Portal", 26, True)
    add_pic(s, labeled("ui-03-worker.png"), Inches(0.35), Inches(0.7), Inches(6.25), Inches(3.2))
    add_pic(s, labeled("ui-05-worker-jobs.png"), Inches(6.8), Inches(0.7), Inches(6.15), Inches(3.2))
    add_bullets(
        s,
        Inches(0.7),
        Inches(4.2),
        Inches(11.8),
        Inches(2.8),
        [
            "Active vs completed job filters with SLA urgency cues",
            "Start and resolve assignments with optional repair cost",
            "Toggle availability so supervisors assign only ready workers",
        ],
        18,
    )

    # 13 Supervisor
    s = slide()
    add_text(s, Inches(0.4), Inches(0.12), Inches(12), Inches(0.35), "Supervisor Dashboard", 24, True)
    add_pic(s, labeled("ui-06-supervisor-assign.png"), Inches(0.25), Inches(0.55), Inches(6.35), Inches(3.25))
    add_pic(s, labeled("ui-07-supervisor-queue.png"), Inches(6.75), Inches(0.55), Inches(6.25), Inches(3.25))
    add_pic(s, labeled("ui-08-supervisor-overdue.png"), Inches(0.25), Inches(3.95), Inches(4.15), Inches(3.25))
    add_pic(s, labeled("ui-09-supervisor-workers.png"), Inches(4.55), Inches(3.95), Inches(4.15), Inches(3.25))
    add_pic(s, labeled("ui-10-supervisor-chronic.png"), Inches(8.85), Inches(3.95), Inches(4.15), Inches(3.25))

    # 14 Admin
    s = slide()
    add_text(s, Inches(0.4), Inches(0.12), Inches(12), Inches(0.35), "Admin Dashboard & Directory", 24, True)
    add_pic(s, labeled("ui-13-admin-reports.png"), Inches(0.25), Inches(0.55), Inches(6.35), Inches(3.2))
    add_pic(s, labeled("ui-14-admin-directory.png"), Inches(6.75), Inches(0.55), Inches(6.25), Inches(3.2))
    add_pic(s, labeled("ui-15-admin-add-student.png"), Inches(0.25), Inches(3.9), Inches(4.15), Inches(3.25))
    add_pic(s, labeled("ui-16-admin-add-worker.png"), Inches(4.55), Inches(3.9), Inches(4.15), Inches(3.25))
    add_pic(s, labeled("ui-17-admin-students-list.png"), Inches(8.85), Inches(3.9), Inches(4.15), Inches(3.25))

    # 15 Tech + demo
    s = slide()
    add_text(s, Inches(0.55), Inches(0.28), Inches(12), Inches(0.5), "Tech Stack & Demo Accounts", 28, True)
    add_card(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(5.7))
    add_text(s, Inches(0.75), Inches(1.3), Inches(5.5), Inches(0.4), "Technology", 18, True, GREEN)
    add_bullets(
        s,
        Inches(0.85),
        Inches(1.9),
        Inches(5.4),
        Inches(4.5),
        [
            "Oracle Database (SQL + PL/SQL)",
            "Node.js + Express + oracledb",
            "HTML / CSS / JavaScript portals",
            "DBeaver / SQL Developer tooling",
            "GitHub version control",
        ],
        17,
    )
    add_card(s, Inches(6.8), Inches(1.05), Inches(6.0), Inches(5.7))
    add_text(s, Inches(7.05), Inches(1.3), Inches(5.5), Inches(0.4), "Demo Login (Password123)", 18, True, PURPLE)
    add_bullets(
        s,
        Inches(7.15),
        Inches(1.9),
        Inches(5.4),
        Inches(4.5),
        [
            "Student — hassan.r@stu.edu",
            "Worker — rashid.i@campus.edu",
            "Supervisor — omar.s@campus.edu",
            "Admin — admin@campus.edu",
            "Local URL — http://localhost:3000",
        ],
        17,
    )

    # 16 GitHub
    s = slide()
    add_text(s, Inches(0.4), Inches(0.15), Inches(12), Inches(0.4), "GitHub Repository & Development History", 22, True)
    add_text(s, Inches(0.45), Inches(0.55), Inches(12), Inches(0.3), GITHUB, 13, False, ACCENT)
    add_pic(s, labeled("github-01-repo.png"), Inches(0.3), Inches(0.95), Inches(6.4), Inches(3.05))
    add_pic(s, labeled("github-02-commits.png"), Inches(6.9), Inches(0.95), Inches(3.5), Inches(6.1))
    add_pic(s, labeled("github-03-contributors.png"), Inches(0.3), Inches(4.15), Inches(6.4), Inches(2.95))

    # 17 Thank you
    s = slide()
    add_text(s, Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.0), "Thank You", 54, True, TEXT, PP_ALIGN.CENTER)
    add_text(
        s,
        Inches(0.8),
        Inches(3.55),
        Inches(11.7),
        Inches(0.45),
        "Questions & Discussion Welcome",
        20,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(4.5),
        Inches(11.7),
        Inches(0.9),
        "Ahnaf Tajwar Sadi  ·  2207104\nCampus Maintenance & Complaint Management System",
        16,
        False,
        ACCENT,
        PP_ALIGN.CENTER,
    )

    prs.save(OUT_PPTX)
    print("PPTX written:", OUT_PPTX)
    print("Slides:", len(prs.slides))


if __name__ == "__main__":
    build()
